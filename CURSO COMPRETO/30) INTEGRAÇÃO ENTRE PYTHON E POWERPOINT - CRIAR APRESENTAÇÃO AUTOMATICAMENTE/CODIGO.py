from pptx import Presentation
from pptx.util import Inches

# Crie uma apresentação vazia
apresentacao = Presentation()

# Crie o primeiro slide
slide1 = apresentacao.slides.add_slide(apresentacao.slide_layouts[0])
titulo1 = slide1.shapes.title
titulo1.text = "Título do Slide 1"
conteudo1 = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(2))
texto1 = conteudo1.text_frame
texto1.text = "Este é o conteúdo do Slide 1."

# Crie o segundo slide
slide2 = apresentacao.slides.add_slide(apresentacao.slide_layouts[1])
titulo2 = slide2.shapes.title
titulo2.text = "Título do Slide 2"
conteudo2 = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(2))
texto2 = conteudo2.text_frame
texto2.text = "Este é o conteúdo do Slide 2."

# Salve a apresentação em um arquivo
apresentacao.save("minha_apresentacao.pptx")
